#!/data/data/com.termux/files/home/.local/bin/python
from __future__ import annotations
import sys
from pathlib import Path
from pptx import Presentation


def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening file: {e}")
        return
    output_file = Path(pptx_path).with_suffix(".txt")
    with open(output_file, "w", encoding="utf-8") as f:
        for slide_num, slide in enumerate(prs.slides, 1):
            f.write(f"\n{'=' * 40}\n")
            f.write(f"Slide {slide_num}\n")
            f.write(f"{'=' * 40}\n\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    f.write(f"{shape.text}\n")
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text)
                            f.write(" | ".join(row_text) + "\n")
                    f.write("\n")
    print(f"Text extracted to: {output_file}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_text.py <input.pptx>")
        sys.exit(1)
    extract_text_from_pptx(sys.argv[1])
